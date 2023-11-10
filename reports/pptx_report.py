from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pandas import DataFrame


def render_pptx_document(
        dataframe: DataFrame,
        report_path: Path,
        template_path: Path
) -> None:
    """
    'profile_id',
    'user_res_id',
    'subscription_res_id'
    'soc_type',
    'source_type',
    'alert_type'
    """
    download_date = datetime.now()
    download_date = download_date.strftime("%d-%m-%Y")
    users_total_count = str(dataframe['profile_id'].unique().size)
    source_total_count = str(dataframe['user_res_id'].unique().size)
    source_vk_count = str(dataframe.loc[dataframe['soc_type'] == 1]['user_res_id'].unique().size)
    source_insta_count = str(dataframe.loc[dataframe['soc_type'] == 4]['user_res_id'].unique().size)
    sub_count = str(dataframe['subscription_res_id'].unique().size)
    alert_sub_count = str(dataframe.groupby(by=['subscription_res_id', 'alert_type'], dropna=True).count().shape[0])
    sub_vk_count = str(dataframe.loc[dataframe['soc_type'] == 1]['subscription_res_id'].unique().size)
    sub_inst_count = str(dataframe.loc[dataframe['soc_type'] == 4]['subscription_res_id'].unique().size)
    alert_type_1_count = str(dataframe.loc[dataframe['alert_type'] == 1]['subscription_res_id'].unique().size)
    alert_type_2_count = str(dataframe.loc[dataframe['alert_type'] == 2]['subscription_res_id'].unique().size)
    alert_type_3_count = str(dataframe.loc[dataframe['alert_type'] == 3]['subscription_res_id'].unique().size)
    alert_type_4_count = str(dataframe.loc[dataframe['alert_type'] == 4]['subscription_res_id'].unique().size)

    replacements = {
        "<<download_date>>": download_date,
        "<<users_total_count>>": users_total_count,
        "<<source_total_count>>": source_total_count,
        "<<source_vk_count>>": source_vk_count,
        "<<source_insta_count>>": source_insta_count,
        "<<sub_count>>": sub_count,
        "<<alert_sub_count>>": alert_sub_count,
        "<<sub_vk_count>>": sub_vk_count,
        "<<sub_inst_count>>": sub_inst_count,
        "<<alert_type_1_count>>": alert_type_1_count,
        "<<alert_type_2_count>>": alert_type_2_count,
        "<<alert_type_3_count>>": alert_type_3_count,
        "<<alert_type_4_count>>": alert_type_4_count
    }

    try:
        presentation = Presentation(template_path)

        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for old_text, new_text in replacements.items():
                                if old_text in run.text:
                                    run.text = run.text.replace(old_text, new_text)

        presentation.save(report_path)

    except FileNotFoundError:
        raise FileNotFoundError(f'По указанному пути {report_path} ничего не найдено!')


__all__ = [
    'render_pptx_document',
]
